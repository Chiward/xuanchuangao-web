import io
import os
from fastapi import UploadFile, HTTPException
import docx
import pptx
import PyPDF2

async def extract_text_from_file(file: UploadFile) -> str:
    """
    Extract text from uploaded file based on its content type or extension.
    """
    filename = file.filename.lower()
    content = await file.read()
    file_stream = io.BytesIO(content)
    
    return _parse_content(file_stream, filename)

def read_text_from_path(file_path: str) -> str:
    """
    Read text from a local file path.
    """
    if not os.path.exists(file_path):
        return ""
        
    filename = os.path.basename(file_path).lower()
    
    # Read file content
    with open(file_path, "rb") as f:
        content = f.read()
    
    file_stream = io.BytesIO(content)
    
    try:
        return _parse_content(file_stream, filename)
    except Exception as e:
        print(f"Error reading local file {file_path}: {e}")
        return ""

def _parse_content(file_stream, filename: str) -> str:
    try:
        if filename.endswith(".docx"):
            return parse_docx(file_stream)
        elif filename.endswith(".pptx"):
            return parse_pptx(file_stream)
        elif filename.endswith(".pdf"):
            return parse_pdf(file_stream)
        elif filename.endswith(".txt"):
            return file_stream.getvalue().decode("utf-8")
        else:
            # Try to read as plain text for other extensions
            try:
                return file_stream.getvalue().decode("utf-8")
            except:
                # If uploaded file, raise HTTP exception
                # If local file, this will propagate up and return empty string
                raise HTTPException(status_code=400, detail="Unsupported file format")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error parsing file: {str(e)}")

def parse_docx(file_stream) -> str:
    doc = docx.Document(file_stream)
    full_text = []
    for para in doc.paragraphs:
        if para.text.strip():
            full_text.append(para.text)
    return "\n".join(full_text)

def parse_pptx(file_stream) -> str:
    prs = pptx.Presentation(file_stream)
    full_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                full_text.append(shape.text)
    return "\n".join(full_text)

def parse_pdf(file_stream) -> str:
    reader = PyPDF2.PdfReader(file_stream)
    full_text = []
    for page in reader.pages:
        text = page.extract_text()
        if text and text.strip():
            full_text.append(text)
    return "\n".join(full_text)
